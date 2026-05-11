"""chunking.py 단위 테스트 — 포맷별 청킹 전략 검증."""
import sys
from pathlib import Path

import pytest

# log-analyzer 루트 디렉터리를 import path에 추가
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import chunking  # noqa: E402
from chunking import (  # noqa: E402
    _find_break_point,
    _merge_short_sections,
    chunk_confluence_page,
    chunk_docx,
    chunk_pdf,
    chunk_pptx,
    chunk_text,
    chunk_xlsx,
    make_ocr_stats,
)


# ── chunk_text ────────────────────────────────────────────────────────────────

def test_chunk_text_empty():
    assert chunk_text("") == []


def test_chunk_text_short():
    chunks = chunk_text("짧은 텍스트", max_chars=1500, overlap=200)
    assert len(chunks) == 1
    assert chunks[0]["text"] == "짧은 텍스트"
    assert chunks[0]["metadata"]["chunk_index"] == 0


def test_chunk_text_5000_chars_yields_4_chunks_with_overlap():
    """5000자(공백 없음) → max=1500, overlap=200 → 4 청크 + overlap 검증.

    공백/줄바꿈이 없으면 _find_break_point가 백트래킹 실패 → 그대로 자름.
    step = max - overlap = 1300
    [0,1500), [1300,2800), [2600,4100), [3900,5000) → 4 청크
    각 인접 청크의 마지막 200자 ⊂ 다음 청크 시작 200자
    """
    text = "가나다라" * 1250  # 5000자, 공백 없음
    assert len(text) == 5000

    chunks = chunk_text(text, max_chars=1500, overlap=200)
    assert len(chunks) == 4
    assert chunks[0]["metadata"]["chunk_index"] == 0
    assert chunks[3]["metadata"]["chunk_index"] == 3

    # 인접 청크 overlap 검증: 청크 N의 끝 200자가 청크 N+1의 시작에 포함
    for i in range(len(chunks) - 1):
        tail = chunks[i]["text"][-200:]
        head = chunks[i + 1]["text"][:200]
        assert tail == head, f"chunk {i}↔{i+1} overlap mismatch"


def test_chunk_text_prefers_paragraph_break():
    """단락(\\n\\n) 경계가 lookback 범위에 있으면 거기서 끊긴다."""
    # 1400자 + \n\n + 200자 → 첫 청크는 1400자 단락 끝에서 끊겨야 함
    para1 = "가" * 1400
    para2 = "나" * 200
    text = para1 + "\n\n" + para2
    chunks = chunk_text(text, max_chars=1500, overlap=200)
    # 첫 청크는 단락 경계에서 끊김 → para1만 포함
    assert chunks[0]["text"] == para1
    # 다음 청크는 para2부터
    assert "나" in chunks[1]["text"]


def test_chunk_text_with_base_metadata():
    chunks = chunk_text("hello world", base_metadata={"source": "manual"})
    assert chunks[0]["metadata"]["source"] == "manual"
    assert chunks[0]["metadata"]["chunk_index"] == 0


def test_chunk_text_invalid_overlap():
    with pytest.raises(ValueError):
        chunk_text("text", max_chars=100, overlap=100)
    with pytest.raises(ValueError):
        chunk_text("text", max_chars=100, overlap=-1)


def test_find_break_point_no_boundary():
    """경계가 없으면 end 그대로."""
    text = "가나다라마바사아자차"
    assert _find_break_point(text, 5, lookback=5) == 5


# ── chunk_confluence_page ─────────────────────────────────────────────────────

def test_chunk_confluence_page_html_with_h2_h3():
    """H2/H3 섹션 분할 + 짧은 섹션 병합 후 모든 내용이 보존되는지 검증."""
    html = """
    <html><body>
      <h1>전체 제목</h1>
      <p>인트로 단락</p>
      <h2>섹션 1</h2>
      <p>섹션 1 본문 첫 단락</p>
      <p>섹션 1 본문 둘째 단락</p>
      <h2>섹션 2</h2>
      <p>섹션 2 본문</p>
      <h3>섹션 2-1</h3>
      <p>하위 섹션 본문</p>
    </body></html>
    """
    chunks = chunk_confluence_page(
        html, page_id="123", page_title="테스트 페이지", space="DOC"
    )
    # 단편 섹션이 많아 병합되므로 청크 수는 1 이상이면 충분
    assert len(chunks) >= 1
    # 모든 청크가 공통 메타 보존
    for c in chunks:
        assert c["metadata"]["source_type"] == "confluence"
        assert c["metadata"]["page_id"] == "123"
        assert c["metadata"]["page_title"] == "테스트 페이지"
        assert c["metadata"]["space"] == "DOC"
    # 모든 본문이 어느 청크에든 포함됨
    all_text = "\n".join(c["text"] for c in chunks)
    assert "인트로 단락" in all_text
    assert "섹션 1 본문 첫 단락" in all_text
    assert "섹션 2 본문" in all_text
    assert "하위 섹션 본문" in all_text


def test_merge_short_sections_forward():
    """짧은 앞 섹션이 다음 섹션에 병합되는지 검증."""
    sections = [
        ("작업 절차", "1단계: 서버 점검"),       # 10자 < 300
        ("주의사항", "점검 전 공지 필수"),         # 9자 < 300
        ("담당자", "홍길동" * 100),               # 300자 이상
    ]
    result = _merge_short_sections(sections, min_chars=300)
    assert len(result) == 1
    h, body = result[0]
    assert h == "작업 절차"           # 체인의 첫 heading 유지
    assert "1단계: 서버 점검" in body
    assert "## 주의사항" in body
    assert "점검 전 공지 필수" in body
    assert "## 담당자" in body
    assert "홍길동" in body


def test_merge_short_sections_last_backward():
    """마지막 섹션이 짧으면 직전 섹션 뒤에 붙는지 검증."""
    long_body = "내용" * 200   # 400자 ≥ 300
    sections = [
        ("본문 섹션", long_body),
        ("짧은 마지막", "한 줄"),   # 3자 < 300
    ]
    result = _merge_short_sections(sections, min_chars=300)
    assert len(result) == 1
    h, body = result[0]
    assert h == "본문 섹션"
    assert long_body in body
    assert "## 짧은 마지막" in body
    assert "한 줄" in body


def test_merge_short_sections_long_sections_unchanged():
    """충분히 긴 섹션은 병합 없이 그대로 유지."""
    long_body = "내용" * 200   # 400자 ≥ 300
    sections = [
        ("섹션 A", long_body),
        ("섹션 B", long_body),
        ("섹션 C", long_body),
    ]
    result = _merge_short_sections(sections, min_chars=300)
    assert len(result) == 3
    assert result[0][0] == "섹션 A"
    assert result[1][0] == "섹션 B"
    assert result[2][0] == "섹션 C"


def test_chunk_confluence_page_short_sections_merged():
    """한 줄짜리 섹션들이 병합되어 포인트 수가 줄어드는지 검증."""
    html = (
        "<body>"
        "<h2>작업 절차</h2><p>1단계: 서버 점검</p>"
        "<h2>주의사항</h2><p>점검 전 공지 필수</p>"
        "<h2>본문</h2><p>" + ("내용" * 200) + "</p>"  # 긴 섹션
        "</body>"
    )
    chunks = chunk_confluence_page(html, page_id="p1", page_title="가이드")
    # 짧은 두 섹션이 긴 섹션에 병합되어 1개 포인트여야 함
    assert len(chunks) == 1
    body = chunks[0]["text"]
    assert "1단계: 서버 점검" in body
    assert "점검 전 공지 필수" in body
    assert "내용" in body


def test_chunk_confluence_page_plain_text():
    """HTML이 아닌 plain text는 chunk_text로 fallback."""
    text = "단순 텍스트 본문입니다."
    chunks = chunk_confluence_page(text, page_id="p1", page_title="title")
    assert len(chunks) == 1
    assert chunks[0]["metadata"]["source_type"] == "confluence"
    assert chunks[0]["metadata"]["page_id"] == "p1"


def test_chunk_confluence_page_long_section_splits():
    """긴 섹션은 sliding window로 분할된다."""
    body = "가" * 3500
    html = f"<body><h2>긴 섹션</h2><p>{body}</p></body>"
    chunks = chunk_confluence_page(html, page_id="x", page_title="t")
    assert len(chunks) >= 2
    # 모두 동일 heading 메타
    for c in chunks:
        assert c["metadata"]["heading"] == "긴 섹션"


def test_chunk_confluence_page_empty():
    assert chunk_confluence_page("", page_id="x", page_title="t") == []


# ── chunk_docx ────────────────────────────────────────────────────────────────

def test_chunk_docx_basic(tmp_path):
    """python-docx로 가짜 docx 생성 → 청킹 결과 검증."""
    docx = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("첫 번째 단락입니다. 한국어 본문 내용.")
    doc.add_paragraph("두 번째 단락. 추가 내용이 들어 있습니다.")
    # 표 추가
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "이름"
    table.rows[0].cells[1].text = "값"
    table.rows[1].cells[0].text = "CPU"
    table.rows[1].cells[1].text = "80%"

    file_path = tmp_path / "sample.docx"
    doc.save(str(file_path))

    chunks = chunk_docx(str(file_path))
    assert len(chunks) >= 1
    combined = "\n".join(c["text"] for c in chunks)
    assert "첫 번째 단락" in combined
    assert "두 번째 단락" in combined
    assert "CPU" in combined and "80%" in combined
    assert chunks[0]["metadata"]["doc_type"] == "docx"
    assert chunks[0]["metadata"]["file_name"] == "sample.docx"


def test_chunk_docx_long_paragraphs_split(tmp_path):
    """1500자 초과 시 여러 청크로 분할."""
    pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    for _ in range(20):
        doc.add_paragraph("가" * 200)  # 단락당 200자, 총 4000자

    file_path = tmp_path / "long.docx"
    doc.save(str(file_path))

    chunks = chunk_docx(str(file_path), max_chars=1500, overlap=200)
    assert len(chunks) >= 2


# ── chunk_pdf (mock) ─────────────────────────────────────────────────────────

def _make_pdf_mocks(monkeypatch, pages_text, pages_images=None):
    """pdfplumber + pypdf mock 헬퍼.

    pages_text: list[str] — 각 페이지 텍스트 (빈 문자열 = 빈 페이지)
    pages_images: list[list[bytes]] | None — 페이지별 이미지 blob 목록
                  None이면 모든 페이지 이미지 없음
    """
    import sys as _sys
    import types

    # ── pdfplumber mock ──────────────────────────────────────────────
    class _MockPdfPage:
        def __init__(self, text):
            self._text = text

        def extract_text(self):
            return self._text

    class _MockPdf:
        def __init__(self, pages):
            self.pages = pages

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

    def _mock_pdfplumber_open(_path):
        return _MockPdf([_MockPdfPage(t) for t in pages_text])

    fake_pdfplumber = types.SimpleNamespace(open=_mock_pdfplumber_open)
    monkeypatch.setitem(_sys.modules, "pdfplumber", fake_pdfplumber)

    # ── pypdf mock ───────────────────────────────────────────────────
    class _MockImageFile:
        def __init__(self, blob):
            self.data = blob

    class _MockPypdfPage:
        def __init__(self, image_blobs):
            self._image_blobs = image_blobs

        @property
        def images(self):
            return [_MockImageFile(b) for b in self._image_blobs]

    class _MockPypdfReader:
        def __init__(self, pages):
            self.pages = pages

    resolved_images = pages_images if pages_images is not None else [[] for _ in pages_text]

    def _mock_pypdf_reader(_path):
        return _MockPypdfReader([_MockPypdfPage(imgs) for imgs in resolved_images])

    fake_pypdf = types.SimpleNamespace(PdfReader=_mock_pypdf_reader)
    monkeypatch.setitem(_sys.modules, "pypdf", fake_pypdf)


def test_chunk_pdf_mocked(monkeypatch):
    """pdfplumber + pypdf를 mock으로 대체 → 페이지별 청킹 검증 (이미지 없음)."""
    _make_pdf_mocks(
        monkeypatch,
        pages_text=[
            "페이지 1 본문입니다.",
            "페이지 2 본문. 좀 더 길게 작성한 내용.",
            "",   # 빈 페이지는 건너뛰어야 함
            "페이지 4 본문.",
        ],
    )

    chunks = chunk_pdf("/fake/path/dummy.pdf")
    assert len(chunks) == 3  # 빈 페이지 제외
    page_nos = [c["metadata"]["page_no"] for c in chunks]
    assert page_nos == [1, 2, 4]
    for c in chunks:
        assert c["metadata"]["doc_type"] == "pdf"
        assert c["metadata"]["file_name"] == "dummy.pdf"
    # chunk_index 전역 누적
    assert [c["metadata"]["chunk_index"] for c in chunks] == [0, 1, 2]


# ── chunk_xlsx ────────────────────────────────────────────────────────────────

def test_chunk_xlsx_basic(tmp_path):
    """openpyxl로 가짜 xlsx 생성 → 시트 단위 청킹 검증."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "메트릭"
    ws1.append(["항목", "값"])
    ws1.append(["CPU", "80%"])
    ws1.append(["MEM", "70%"])

    ws2 = wb.create_sheet(title="알림")
    ws2.append(["시각", "메시지"])
    ws2.append(["10:00", "고부하"])

    file_path = tmp_path / "report.xlsx"
    wb.save(str(file_path))

    chunks = chunk_xlsx(str(file_path))
    assert len(chunks) == 2  # 시트 = 청크
    sheet_names = [c["metadata"]["sheet_name"] for c in chunks]
    assert sheet_names == ["메트릭", "알림"]
    for c in chunks:
        assert c["metadata"]["doc_type"] == "xlsx"
        assert c["metadata"]["file_name"] == "report.xlsx"
        assert "|" in c["text"]  # markdown 표 형식
    # 메트릭 시트 본문에 데이터 포함
    metric_chunk = next(c for c in chunks if c["metadata"]["sheet_name"] == "메트릭")
    assert "CPU" in metric_chunk["text"]
    assert "80%" in metric_chunk["text"]


def test_chunk_xlsx_under_embed_limit_not_split(tmp_path):
    """1500자 초과이지만 _EMBED_TOKEN_LIMIT_CHARS(6500자) 이내 시트: 단일 청크 유지."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "medium"
    ws.append(["col1", "col2"])
    # 행당 ~30자 × 80행 ≈ 2400자 → 6500 이내
    for i in range(80):
        ws.append([f"value-{i}-가나다라마", f"val-{i}"])

    file_path = tmp_path / "medium.xlsx"
    wb.save(str(file_path))

    chunks = chunk_xlsx(str(file_path))
    assert len(chunks) == 1  # 1500 초과지만 6500 이내 → 분할 안 됨
    assert len(chunks[0]["text"]) > 1500
    # 단일 청크 시에도 sub_chunk_index=0, total_sub_chunks=1 포함
    assert chunks[0]["metadata"]["sub_chunk_index"] == 0
    assert chunks[0]["metadata"]["total_sub_chunks"] == 1


# ── chunk_pptx ────────────────────────────────────────────────────────────────

def test_chunk_pptx_basic(tmp_path):
    """python-pptx로 가짜 pptx 생성 → 슬라이드 단위 청킹 검증."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    # 슬라이드 1 — 제목 + 본문
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "첫 슬라이드"
    body_placeholder = slide1.placeholders[1]
    body_placeholder.text = "첫 슬라이드 본문 내용"
    # 발표자 노트 추가
    slide1.notes_slide.notes_text_frame.text = "발표자 메모입니다"

    # 슬라이드 2 — 제목만
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "두 번째 슬라이드"

    file_path = tmp_path / "deck.pptx"
    prs.save(str(file_path))

    chunks = chunk_pptx(str(file_path))
    assert len(chunks) == 2
    # 슬라이드 1
    assert chunks[0]["metadata"]["slide_no"] == 1
    assert chunks[0]["metadata"]["slide_title"] == "첫 슬라이드"
    assert chunks[0]["metadata"]["doc_type"] == "pptx"
    assert chunks[0]["metadata"]["file_name"] == "deck.pptx"
    assert "첫 슬라이드 본문 내용" in chunks[0]["text"]
    assert "발표자 메모" in chunks[0]["text"]
    # 슬라이드 2
    assert chunks[1]["metadata"]["slide_no"] == 2
    assert chunks[1]["metadata"]["slide_title"] == "두 번째 슬라이드"


# ── _EMBED_TOKEN_LIMIT_CHARS 상수 ────────────────────────────────────────────

def test_embed_token_limit_constant_exists():
    """_EMBED_TOKEN_LIMIT_CHARS 상수가 6500으로 정의되어 있는지 검증."""
    assert chunking._EMBED_TOKEN_LIMIT_CHARS == 6500


# ── chunk_xlsx 분할 (임계값 초과) ─────────────────────────────────────────────

def _make_xlsx_with_rows(tmp_path, sheet_name: str, num_rows: int, row_len: int = 30):
    """지정한 행 수의 xlsx 파일 생성 헬퍼. 행당 row_len자 내외 데이터."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    ws.append(["헤더1", "헤더2"])
    for i in range(num_rows):
        # 행당 약 row_len 자 채우기
        pad = "가" * max(1, row_len // 2 - len(str(i)))
        ws.append([f"{pad}{i}", f"{pad}{i}"])
    file_path = tmp_path / f"{sheet_name}.xlsx"
    wb.save(str(file_path))
    return str(file_path)


def test_chunk_xlsx_at_limit_not_split(tmp_path):
    """6500자 이내 시트 — 분할 안 됨, oversize_count=0."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "small"
    ws.append(["col"])
    # 20행 × ~30자 = ~600자 → 6500 이내
    for i in range(20):
        ws.append([f"val-{i}-" + "나" * 10])

    file_path = tmp_path / "small.xlsx"
    wb.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_xlsx(str(file_path), stats=stats)

    assert len(chunks) == 1
    assert chunks[0]["metadata"]["sub_chunk_index"] == 0
    assert chunks[0]["metadata"]["total_sub_chunks"] == 1
    assert stats["oversize_count"] == 0


def test_chunk_xlsx_oversize_splits_with_header(tmp_path):
    """7000자 시트 — 2개 sub-chunk, 헤더가 양쪽에 보존됨."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "big"
    ws.append(["응답시간", "에러코드"])
    # 행당 ~40자 × 180행 ≈ 7200자 → 6500 초과
    for i in range(180):
        ws.append([f"응답시간-{i:04d}-{'가'*10}", f"에러코드-{i:04d}"])

    file_path = tmp_path / "big.xlsx"
    wb.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_xlsx(str(file_path), stats=stats)

    # 분할 발생
    assert len(chunks) >= 2
    # oversize_count 누적
    assert stats["oversize_count"] == 1
    # 모든 sub-chunk에 헤더 보존 (키워드 "응답시간" 검색 유지)
    for c in chunks:
        assert "응답시간" in c["text"]
        assert "에러코드" in c["text"]
    # sub_chunk_index/total_sub_chunks 메타 일관성
    total = chunks[0]["metadata"]["total_sub_chunks"]
    assert total == len(chunks)
    for idx, c in enumerate(chunks):
        assert c["metadata"]["sub_chunk_index"] == idx
        assert c["metadata"]["total_sub_chunks"] == total
    # 각 sub-chunk는 임계값 이내
    for c in chunks:
        assert len(c["text"]) <= chunking._EMBED_TOKEN_LIMIT_CHARS


def test_chunk_xlsx_very_large_sheet_multiple_subchunks(tmp_path):
    """30000자 시트 — 5개 이상 sub-chunk, 헤더 모두 보존."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "huge"
    ws.append(["항목명", "상세내용"])
    # 행당 ~60자 × 500행 ≈ 30000자
    for i in range(500):
        ws.append([f"항목-{i:04d}-{'나'*15}", f"상세-{i:04d}-{'다'*15}"])

    file_path = tmp_path / "huge.xlsx"
    wb.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_xlsx(str(file_path), stats=stats)

    assert len(chunks) >= 5
    assert stats["oversize_count"] == 1
    for c in chunks:
        assert "항목명" in c["text"]
        assert "상세내용" in c["text"]
        assert len(c["text"]) <= chunking._EMBED_TOKEN_LIMIT_CHARS


def test_chunk_xlsx_stats_none_no_error(tmp_path):
    """stats=None (기본) 시 oversize 분할해도 예외 없음 (BC 유지)."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "test"
    ws.append(["col"])
    for i in range(180):
        ws.append([f"value-{i}-{'가'*20}"])

    file_path = tmp_path / "nostat.xlsx"
    wb.save(str(file_path))

    # stats 미전달 — 예외 없이 정상 동작
    chunks = chunk_xlsx(str(file_path))
    assert len(chunks) >= 1


# ── chunk_pptx 분할 (임계값 초과) ────────────────────────────────────────────

def test_chunk_pptx_small_slide_single_chunk(tmp_path):
    """작은 슬라이드 — 단일 청크, sub_chunk_index=0."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "제목"
    slide.placeholders[1].text = "짧은 본문"

    file_path = tmp_path / "small.pptx"
    prs.save(str(file_path))

    chunks = chunk_pptx(str(file_path))
    assert len(chunks) == 1
    assert chunks[0]["metadata"]["sub_chunk_index"] == 0
    assert chunks[0]["metadata"]["total_sub_chunks"] == 1


def test_chunk_pptx_body_notes_separated(tmp_path):
    """body=5000자 + notes=3000자 → body, notes 분리, title 양쪽 prepend."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "장애 분석 보고서"
    slide.placeholders[1].text = "가" * 5000
    slide.notes_slide.notes_text_frame.text = "나" * 3000

    file_path = tmp_path / "split.pptx"
    prs.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_pptx(str(file_path), stats=stats)

    # body + notes 분리 → 2개
    assert len(chunks) == 2
    assert stats["oversize_count"] == 1
    # title이 양쪽에 prepend
    for c in chunks:
        assert "장애 분석 보고서" in c["text"]
    # 각 sub-chunk 임계값 이내
    for c in chunks:
        assert len(c["text"]) <= chunking._EMBED_TOKEN_LIMIT_CHARS
    # sub_chunk_index/total_sub_chunks 일관성
    assert chunks[0]["metadata"]["sub_chunk_index"] == 0
    assert chunks[0]["metadata"]["total_sub_chunks"] == 2
    assert chunks[1]["metadata"]["sub_chunk_index"] == 1
    assert chunks[1]["metadata"]["total_sub_chunks"] == 2
    # body 청크에 가 포함, notes 청크에 나 포함
    assert "가" * 100 in chunks[0]["text"]
    assert "나" * 100 in chunks[1]["text"]


def test_chunk_pptx_body_oversize_further_split(tmp_path):
    """body=10000자 (단독 초과) → chunk_text로 추가 분할."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "초대형 슬라이드"
    slide.placeholders[1].text = "가" * 10000

    file_path = tmp_path / "huge_slide.pptx"
    prs.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_pptx(str(file_path), stats=stats)

    # body 단독 10000자 → 여러 sub-chunk
    assert len(chunks) >= 2
    assert stats["oversize_count"] == 1
    for c in chunks:
        assert len(c["text"]) <= chunking._EMBED_TOKEN_LIMIT_CHARS
        # title_prefix가 모든 sub-chunk에 보존되어야 한다
        assert "초대형 슬라이드" in c["text"]


def test_chunk_pptx_oversize_count_per_slide(tmp_path):
    """여러 슬라이드 중 초과 슬라이드 수만큼 oversize_count 누적."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()

    # 슬라이드 1: 정상 크기
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "정상 슬라이드"
    s1.placeholders[1].text = "짧은 내용"

    # 슬라이드 2: 초과 크기
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "큰 슬라이드"
    s2.placeholders[1].text = "가" * 5000
    s2.notes_slide.notes_text_frame.text = "나" * 3000

    # 슬라이드 3: 초과 크기
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "또다른 큰 슬라이드"
    s3.placeholders[1].text = "다" * 5000
    s3.notes_slide.notes_text_frame.text = "라" * 3000

    file_path = tmp_path / "multi.pptx"
    prs.save(str(file_path))

    stats = make_ocr_stats()
    chunks = chunk_pptx(str(file_path), stats=stats)

    # 슬라이드 2, 3 초과 → oversize_count=2
    assert stats["oversize_count"] == 2
    # 전체 청크 수: 슬라이드1(1) + 슬라이드2(2) + 슬라이드3(2) = 5
    assert len(chunks) == 5


# ── chunk_pdf sliding window 회귀 확인 ───────────────────────────────────────

def _make_pdf_mocks_local(monkeypatch, pages_text, pages_images=None):
    """pdfplumber + pypdf mock 헬퍼 (로컬 재정의, test_chunking.py 전용)."""
    import sys as _sys
    import types

    class _MockPdfPage:
        def __init__(self, text):
            self._text = text

        def extract_text(self):
            return self._text

    class _MockPdf:
        def __init__(self, pages):
            self.pages = pages

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

    monkeypatch.setitem(
        _sys.modules,
        "pdfplumber",
        types.SimpleNamespace(open=lambda _: _MockPdf([_MockPdfPage(t) for t in pages_text])),
    )

    class _MockImageFile:
        def __init__(self, blob):
            self.data = blob

    class _MockPypdfPage:
        def __init__(self, blobs):
            self._blobs = blobs

        @property
        def images(self):
            return [_MockImageFile(b) for b in self._blobs]

    class _MockPypdfReader:
        def __init__(self, pages):
            self.pages = pages

    resolved = pages_images if pages_images is not None else [[] for _ in pages_text]
    monkeypatch.setitem(
        _sys.modules,
        "pypdf",
        types.SimpleNamespace(
            PdfReader=lambda _path: _MockPypdfReader(
                [_MockPypdfPage(imgs) for imgs in resolved]
            )
        ),
    )


def test_chunk_pdf_large_page_sliding_window(monkeypatch):
    """PDF 페이지가 6500자 초과해도 sliding window(chunk_text)로 정상 분할됨 (회귀 확인)."""
    _make_pdf_mocks_local(
        monkeypatch,
        pages_text=["가" * 8000],  # 8000자 단일 페이지 → chunk_text로 분할
    )

    chunks = chunk_pdf("/fake/large_page.pdf")
    assert len(chunks) >= 2
    for c in chunks:
        assert c["metadata"]["page_no"] == 1
        assert c["metadata"]["doc_type"] == "pdf"
    # chunk_index 전역 누적
    assert [c["metadata"]["chunk_index"] for c in chunks] == list(range(len(chunks)))
