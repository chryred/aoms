"""문서 포맷별 청킹 전략 (한국어 RAG 최적화).

청크 크기: 1500자 (한국어 ≈ 800~1000 토큰, bge-m3 8192 한도 내 안전 마진)
오버랩: 200자 (의미 연결 보존)

각 함수는 list[dict] 반환:
  [{"text": "...", "metadata": {"chunk_index": 0, "source_type": "...", ...}}]

설계 원칙:
- 순수 텍스트 sliding window는 ``chunk_text``를 베이스로 모든 포맷이 재사용
- 한국어 청크 경계는 단어/조사 중간을 피하기 위해 단락(\\n\\n) → 줄바꿈(\\n) → 공백 순으로 백트래킹
- xlsx/pptx는 의미 단위(시트/슬라이드)가 곧 청크. bge-m3 8192 토큰 한도(≈ 한국어 6500자)
  초과 시에만 분할 — 헤더/타이틀을 모든 sub-chunk에 복사해 의미 보존
- vector_client.py 등 기존 모듈은 수정하지 않음 (이 모듈은 독립 유틸리티)
"""

from __future__ import annotations

import logging
import os
from typing import Any

logger = logging.getLogger(__name__)

_EMBED_TOKEN_LIMIT_CHARS = 6500  # bge-m3 8192 토큰 ≈ 한국어 6500자 (1.2x 안전 마진)


# ── 베이스: 텍스트 sliding window ──────────────────────────────────────────────

def _find_break_point(text: str, end: int, lookback: int) -> int:
    """``end`` 위치 기준으로 단락→줄바꿈→공백 순으로 거꾸로 탐색해 끊기 좋은 지점 반환.

    ``lookback`` 범위 내에서 적절한 경계를 찾지 못하면 ``end``를 그대로 돌려준다.
    한국어 청크가 조사/단어 중간에서 끊기는 것을 막기 위한 보조 함수.
    """
    if end >= len(text):
        return len(text)
    window_start = max(0, end - lookback)
    # 단락 경계 우선
    para_break = text.rfind("\n\n", window_start, end)
    if para_break != -1 and para_break > window_start:
        return para_break + 2
    # 줄바꿈
    newline = text.rfind("\n", window_start, end)
    if newline != -1 and newline > window_start:
        return newline + 1
    # 공백
    space = text.rfind(" ", window_start, end)
    if space != -1 and space > window_start:
        return space + 1
    # 적절한 경계 없으면 그대로 자름
    return end


def chunk_text(
    text: str,
    max_chars: int = 1500,
    overlap: int = 200,
    base_metadata: dict | None = None,
) -> list[dict]:
    """순수 텍스트 sliding window 청킹 (베이스 함수).

    - ``max_chars``: 청크 최대 길이 (한국어 1500자 ≈ 800~1000 토큰 권장)
    - ``overlap``: 인접 청크 간 중첩 길이 (의미 연결 보존)
    - ``base_metadata``: 모든 청크에 공통으로 박을 메타데이터 (선택)

    경계가 단어 중간이면 ``_find_break_point``로 단락/줄바꿈/공백 위치까지 백트래킹한다.
    """
    if not text:
        return []
    if max_chars <= 0:
        raise ValueError("max_chars must be > 0")
    if overlap < 0 or overlap >= max_chars:
        raise ValueError("overlap must be in [0, max_chars)")

    base_meta = dict(base_metadata) if base_metadata else {}
    chunks: list[dict] = []
    n = len(text)
    start = 0
    chunk_index = 0

    while start < n:
        tentative_end = min(start + max_chars, n)
        if tentative_end < n:
            end = _find_break_point(text, tentative_end, lookback=overlap)
            # 백트래킹이 너무 짧게 끊으면 그대로 사용
            if end <= start:
                end = tentative_end
        else:
            end = n

        piece = text[start:end].strip()
        if piece:
            meta = dict(base_meta)
            meta["chunk_index"] = chunk_index
            chunks.append({"text": piece, "metadata": meta})
            chunk_index += 1

        if end >= n:
            break
        # 다음 시작점: end - overlap (단, 진행을 보장)
        next_start = end - overlap
        if next_start <= start:
            next_start = start + 1
        start = next_start

    return chunks


# ── 이미지 OCR 공통 헬퍼 ─────────────────────────────────────────────────────

_OCR_MIN_LEN = 10           # 의미있는 텍스트 최소 길이
_OCR_VALID_RATIO = 0.7      # 정상 문자 비율 임계 (가비지 필터)


def _is_meaningful_ocr(text: str) -> bool:
    """OCR 결과가 노이즈인지 판정. 한글/영문/숫자/공백/일반 구두점 비율로 추정."""
    if len(text) < _OCR_MIN_LEN:
        return False
    valid = sum(1 for c in text if c.isalnum() or c in " \n\t.,()[]/-:;%")
    return (valid / len(text)) >= _OCR_VALID_RATIO


def make_ocr_stats() -> dict:
    """OCR 통계 카운터 초기화. chunk_* 호출자가 수집에 사용."""
    return {
        "ocr_attempted": 0,
        "ocr_succeeded": 0,
        "ocr_noise_filtered": 0,
        "ocr_failed": 0,
        "oversize_count": 0,
    }


def _ocr_image_blob(blob: bytes, stats: dict | None = None) -> str:
    """이미지 blob → Tesseract OCR 텍스트. 실패/노이즈는 빈 문자열.

    lang=kor+eng, timeout=3s. import는 함수 내부에서 lazy로 호출해
    Tesseract/Pillow 미설치 환경(테스트 등)에서 모듈 import 자체는 깨지지 않도록 한다.

    stats: make_ocr_stats() 반환 dict를 전달하면 ocr_attempted/ocr_succeeded/
           ocr_noise_filtered/ocr_failed 카운터를 누적한다. None이면 통계 생략(BC 유지).
    """
    if stats is not None:
        stats["ocr_attempted"] += 1
    try:
        from PIL import Image
        import pytesseract
        import io
        img = Image.open(io.BytesIO(blob))
        raw = pytesseract.image_to_string(img, lang="kor+eng", timeout=3).strip()
        if _is_meaningful_ocr(raw):
            if stats is not None:
                stats["ocr_succeeded"] += 1
            return raw
        logger.debug("OCR 노이즈 폐기: %d자 (의미없음)", len(raw))
        if stats is not None:
            stats["ocr_noise_filtered"] += 1
        return ""
    except Exception as e:
        logger.warning("OCR exception: %s: %s", type(e).__name__, e)
        if stats is not None:
            stats["ocr_failed"] += 1
        return ""


# ── Confluence 페이지 (HTML or 텍스트) ─────────────────────────────────────────

_MIN_SECTION_CHARS = 300  # 이 미만 섹션은 인접 섹션과 병합


def _merge_short_sections(
    sections: list[tuple[str, str]],
    min_chars: int = _MIN_SECTION_CHARS,
) -> list[tuple[str, str]]:
    """300자 미만 섹션을 인접 섹션과 병합해 단독 포인트 생성 방지.

    짧은 섹션은 다음 섹션 body 앞에 '## heading\\nbody' 형태로 이어붙인다.
    첫 heading은 병합 체인 내내 유지 — 최종 포인트의 metadata.heading이 됨.
    마지막 섹션이 짧으면 직전 확정 섹션 뒤에 붙인다 (뒤로 병합).
    """
    if not sections:
        return sections

    result: list[tuple[str, str]] = []
    acc_heading, acc_body = sections[0]

    for heading, body in sections[1:]:
        if len(acc_body) < min_chars:
            # acc 섹션이 짧으면 현재 섹션을 acc에 흡수
            if heading:
                acc_body = f"{acc_body}\n\n## {heading}\n{body}".strip()
            else:
                acc_body = f"{acc_body}\n\n{body}".strip()
            # acc_heading은 체인의 첫 heading 유지
        else:
            result.append((acc_heading, acc_body))
            acc_heading, acc_body = heading, body

    # 마지막 acc 처리
    if acc_body:
        if result and len(acc_body) < min_chars:
            # 마지막 섹션이 짧으면 직전 확정 섹션 뒤에 붙임
            prev_h, prev_b = result.pop()
            suffix = f"## {acc_heading}\n{acc_body}" if acc_heading else acc_body
            result.append((prev_h, f"{prev_b}\n\n{suffix}".strip()))
        else:
            result.append((acc_heading, acc_body))

    return result


_TABLE_CELL = {"td", "th"}
_TEXT_BLOCK = {"p", "li", "pre", "blockquote", "h1", "h4", "h5", "h6"}
_SECTION_BREAK = {"h2", "h3"}
_SKIP_TAGS = {"script", "style"}
_IMAGE_TAGS = {"ac:image", "img"}


def _walk_html(root) -> list[tuple[str, str]]:
    """HTML 트리를 children 단위로 재귀 탐색해 (kind, text) 목록 반환.

    - ("break", heading): h2/h3 섹션 경계
    - ("text",  content): 단락/리스트/표 행 텍스트

    root.descendants 대신 children 재귀를 사용해
    <th><h2>...</h2></th> 구조에서 h2가 섹션 경계로 오인되는 버그와
    td/th 중복 텍스트 추출 버그를 함께 수정한다.
    table은 행 단위로 '|' 결합해 열 관계를 보존한다.
    """
    results: list[tuple[str, str]] = []
    for elem in root.children:
        name = getattr(elem, "name", None)
        if name is None:
            continue
        if name in _SKIP_TAGS:
            continue
        if name in _SECTION_BREAK:
            results.append(("break", elem.get_text(strip=True)))
        elif name == "table":
            for tr in elem.find_all("tr"):
                cells = [
                    c.get_text(separator=" ", strip=True)
                    for c in tr.find_all(["td", "th"])
                ]
                row = " | ".join(c for c in cells if c)
                if row:
                    results.append(("text", row))
        elif name in _IMAGE_TAGS:
            alt = (elem.get("alt") or elem.get("ac:alt") or "").strip()
            attachment = elem.find("ri:attachment")
            filename = attachment.get("ri:filename", "") if attachment else ""
            label = alt or filename or "이미지"
            results.append(("text", f"[이미지: {label}]"))
            # 자식 재귀 안 함 (이미지 내부에 텍스트 노드 없음)
        elif name in _TEXT_BLOCK:
            txt = elem.get_text(separator=" ", strip=True)
            if txt:
                results.append(("text", txt))
        else:
            results.extend(_walk_html(elem))
    return results


def _looks_like_html(text: str) -> bool:
    snippet = text[:512].lower()
    return "<" in snippet and (">" in snippet) and any(
        tag in snippet for tag in ("<p", "<div", "<h1", "<h2", "<h3", "<span", "<br", "<ul", "<ol", "<table")
    )


def chunk_confluence_page(
    content: str,
    page_id: str,
    page_title: str,
    space: str = "",
    **extra_meta: Any,
) -> list[dict]:
    """Confluence 페이지: H2/H3 heading 우선 분할 → 큰 섹션은 sliding window.

    - HTML이면 BeautifulSoup으로 파싱 후 H2/H3 경계로 섹션 분할
    - plain text면 ``chunk_text``로 바로 분할
    - 각 섹션이 1500자를 넘으면 ``chunk_text``를 다시 적용
    - 메타에 heading(있으면), page_id, page_title, space, source_type='confluence' 보존
    """
    if not content:
        return []

    base_meta: dict[str, Any] = {
        "source_type": "confluence",
        "page_id": page_id,
        "page_title": page_title,
    }
    if space:
        base_meta["space"] = space
    for k, v in extra_meta.items():
        base_meta[k] = v

    if not _looks_like_html(content):
        return chunk_text(content, base_metadata=base_meta)

    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")

    # 본문 흐름 순서대로 H2/H3을 경계로 섹션을 끊는다.
    sections: list[tuple[str, str]] = []  # (heading, text)
    current_heading = ""
    current_buffer: list[str] = []

    # body 또는 root 직속 자식들을 순회해 헤딩 기준으로 분할
    root = soup.body if soup.body else soup
    for kind, text in _walk_html(root):
        if kind == "break":
            buf_text = "\n".join(s for s in current_buffer if s).strip()
            if buf_text:
                sections.append((current_heading, buf_text))
            current_heading = text
            current_buffer = []
        else:
            current_buffer.append(text)

    # 마지막 섹션 마무리
    tail = "\n".join(s for s in current_buffer if s).strip()
    if tail:
        sections.append((current_heading, tail))

    # 헤딩이 하나도 없었던 경우(=descendants 순회로 잡히는 게 없음): 전체 텍스트를 하나로
    if not sections:
        plain = soup.get_text(separator="\n", strip=True)
        return chunk_text(plain, base_metadata=base_meta)

    # 300자 미만 짧은 섹션을 인접 섹션과 병합 — 한 줄짜리 단독 포인트 방지
    sections = _merge_short_sections(sections)

    chunks: list[dict] = []
    chunk_index = 0
    for heading, body in sections:
        if not body:
            continue
        section_meta = dict(base_meta)
        if heading:
            section_meta["heading"] = heading
        # 섹션이 작으면 1 청크, 크면 sliding window
        if len(body) <= 1500:
            section_meta["chunk_index"] = chunk_index
            chunks.append({"text": body, "metadata": section_meta})
            chunk_index += 1
        else:
            sub_chunks = chunk_text(body, base_metadata=section_meta)
            for sc in sub_chunks:
                # chunk_text는 chunk_index를 0부터 부여 → 전역 인덱스로 재할당
                sc["metadata"]["chunk_index"] = chunk_index
                chunks.append(sc)
                chunk_index += 1

    return chunks


# ── DOCX ─────────────────────────────────────────────────────────────────────

def chunk_docx(
    file_path: str,
    max_chars: int = 1500,
    overlap: int = 200,
    stats: dict | None = None,
) -> list[dict]:
    """DOCX 파일: paragraphs 합쳐서 sliding window 청킹.

    - paragraphs와 tables(행 단위)에서 텍스트 추출
    - 단락 사이는 \\n\\n으로 결합 → ``chunk_text``의 단락 경계 백트래킹과 결합
    - metadata: ``{file_name, doc_type: "docx"}``
    - stats: make_ocr_stats() dict를 전달하면 OCR 카운터 누적 (None이면 생략, BC 유지)
    """
    from docx import Document

    doc = Document(file_path)
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            row_text = " | ".join(c for c in row_cells if c)
            if row_text:
                parts.append(row_text)

    # inline image OCR (related_parts에서 image content_type만 추출)
    for rel_id, rel in doc.part.related_parts.items():
        content_type = getattr(rel, "content_type", "") or ""
        if not content_type.startswith("image/"):
            continue
        blob = getattr(rel, "blob", None)
        if not blob:
            continue
        ocr_text = _ocr_image_blob(blob, stats=stats)
        if ocr_text:
            parts.append(f"[이미지: {ocr_text}]")

    full_text = "\n\n".join(parts)
    base_meta = {
        "source_type": "docx",
        "doc_type": "docx",
        "file_name": os.path.basename(file_path),
    }
    return chunk_text(full_text, max_chars=max_chars, overlap=overlap, base_metadata=base_meta)


# ── PDF ──────────────────────────────────────────────────────────────────────

_PDF_MAX_IMAGES_PER_PAGE = 10  # 페이지당 OCR 시도 상한 (표지·차트 PDF 타임아웃 방지)


def chunk_pdf(
    file_path: str,
    max_chars: int = 1500,
    overlap: int = 200,
    stats: dict | None = None,
) -> list[dict]:
    """PDF 파일: 페이지별 텍스트 + inline image OCR 후 sliding window 청킹.

    - pdfplumber로 페이지 단위 텍스트 추출
    - pypdf로 페이지별 embedded image bytes 추출 → ``_ocr_image_blob`` → ``[이미지: ...]`` 마커
      (pypdf가 DCTDecode/FlateDecode/CCITTFaxDecode 등 인코딩을 PIL 호환 bytes로 정규화)
    - 텍스트와 이미지 마커를 ``\\n\\n`` 결합 후 ``chunk_text``로 분할
    - OCR 실패(예외·노이즈)는 silent — 청킹 자체가 깨지지 않도록 보장
    - 페이지당 이미지 OCR 상한: ``_PDF_MAX_IMAGES_PER_PAGE`` (기본 10장)
    - metadata: ``{file_name, doc_type: "pdf", page_no}``
    - 청크 인덱스는 문서 전역으로 누적
    - stats: make_ocr_stats() dict를 전달하면 OCR 카운터 누적 (None이면 생략, BC 유지)
      상한 초과로 건너뛴 이미지(skipped_cap)는 ocr_attempted에 포함하지 않음
    """
    import pdfplumber
    import pypdf

    file_name = os.path.basename(file_path)
    chunks: list[dict] = []
    chunk_index = 0

    # pypdf reader: 이미지 추출 전용 (read-only, pdfplumber와 독립)
    try:
        pypdf_reader = pypdf.PdfReader(file_path)
    except Exception as exc:
        logger.warning("chunk_pdf: pypdf failed to open %s (%s) — image OCR skipped", file_name, exc)
        pypdf_reader = None

    with pdfplumber.open(file_path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            page_text = page_text.strip()

            # ── inline image OCR ──────────────────────────────────────────
            image_markers: list[str] = []
            if pypdf_reader is not None:
                pypdf_page_idx = page_no - 1
                if pypdf_page_idx < len(pypdf_reader.pages):
                    pypdf_page = pypdf_reader.pages[pypdf_page_idx]
                    try:
                        page_images = list(pypdf_page.images)
                    except Exception as exc:
                        logger.debug(
                            "chunk_pdf: page=%d image list failed (%s)", page_no, exc
                        )
                        page_images = []

                    page_ocr_attempted = 0
                    page_ocr_success = 0
                    for img_file in page_images[:_PDF_MAX_IMAGES_PER_PAGE]:
                        try:
                            blob = img_file.data
                        except Exception as exc:
                            logger.warning(
                                "chunk_pdf: page=%d image blob extraction failed (%s)",
                                page_no,
                                exc,
                            )
                            continue
                        if not blob:
                            continue
                        page_ocr_attempted += 1
                        ocr_text = _ocr_image_blob(blob, stats=stats)
                        if ocr_text:
                            image_markers.append(f"[이미지: {ocr_text}]")
                            page_ocr_success += 1

                    skipped = max(0, len(page_images) - _PDF_MAX_IMAGES_PER_PAGE)
                    logger.debug(
                        "chunk_pdf: page=%d images=%d ocr_attempted=%d ocr_success=%d skipped_cap=%d",
                        page_no,
                        len(page_images),
                        page_ocr_attempted,
                        page_ocr_success,
                        skipped,
                    )
            # ─────────────────────────────────────────────────────────────

            # 텍스트 + 이미지 마커 결합
            parts: list[str] = []
            if page_text:
                parts.append(page_text)
            parts.extend(image_markers)

            combined = "\n\n".join(parts).strip()
            if not combined:
                continue

            page_meta = {
                "source_type": "pdf",
                "doc_type": "pdf",
                "file_name": file_name,
                "page_no": page_no,
            }
            sub = chunk_text(combined, max_chars=max_chars, overlap=overlap, base_metadata=page_meta)
            for c in sub:
                c["metadata"]["chunk_index"] = chunk_index
                chunks.append(c)
                chunk_index += 1

    return chunks


# ── XLSX ─────────────────────────────────────────────────────────────────────

def _sheet_to_markdown(ws) -> str:
    """openpyxl Worksheet → markdown 표 형태 텍스트.

    - 첫 행을 헤더로 가정. 첫 행이 비어 있으면 그냥 cell 텍스트만 join.
    - 빈 행/완전히 비어 있는 시트는 ""를 돌려준다.
    """
    rows = list(ws.iter_rows(values_only=True))
    # 빈 셀만 있는 행 제거
    rows = [tuple("" if c is None else str(c) for c in row) for row in rows]
    rows = [row for row in rows if any(c.strip() for c in row)]
    if not rows:
        return ""

    # 열 폭 정규화 (가장 긴 행 기준)
    max_cols = max(len(r) for r in rows)
    rows = [r + ("",) * (max_cols - len(r)) for r in rows]

    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []

    lines: list[str] = []
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def chunk_xlsx(file_path: str, stats: dict | None = None) -> list[dict]:
    """Excel 파일: 시트별 1 chunk (시트 = 표 단위 의미 묶음).

    - openpyxl 사용, 시트 전체를 markdown 표 형태로 변환
    - bge-m3 한도(``_EMBED_TOKEN_LIMIT_CHARS`` = 6500자) 이내 시트: 단일 청크 (BC 유지)
    - 초과 시: 헤더 행을 모든 sub-chunk에 복사하며 데이터 행 단위 분할
    - metadata: ``{file_name, sheet_name, doc_type: "xlsx", chunk_index, sub_chunk_index, total_sub_chunks}``
      (단일 청크 시에도 sub_chunk_index=0, total_sub_chunks=1 포함 — 일관된 메타 형태)
    - stats: make_ocr_stats() dict를 전달하면 oversize_count 누적 (None이면 생략, BC 유지)
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    file_name = os.path.basename(file_path)

    chunks: list[dict] = []
    chunk_index = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_text = _sheet_to_markdown(ws)
        if not sheet_text:
            continue
        # 시트명을 본문 맨 앞에 붙여 검색에 활용 가능하도록
        body = f"# {sheet_name}\n\n{sheet_text}"

        if len(body) <= _EMBED_TOKEN_LIMIT_CHARS:
            # ── 단일 청크 (임계값 이내) ────────────────────────────────────
            meta = {
                "source_type": "xlsx",
                "doc_type": "xlsx",
                "file_name": file_name,
                "sheet_name": sheet_name,
                "chunk_index": chunk_index,
                "sub_chunk_index": 0,
                "total_sub_chunks": 1,
            }
            chunks.append({"text": body, "metadata": meta})
            chunk_index += 1
        else:
            # ── 임계값 초과: 헤더 보존하며 데이터 행 분할 ─────────────────
            if stats is not None:
                stats["oversize_count"] += 1

            lines = sheet_text.splitlines()
            # lines[0]: | header | 행, lines[1]: | --- | 행, lines[2:]: 데이터 행
            if len(lines) >= 2:
                header_lines = lines[0] + "\n" + lines[1]
                data_lines = lines[2:]
            else:
                # 구분선 없는 드문 경우 — 첫 행만 헤더로 취급
                header_lines = lines[0]
                data_lines = lines[1:]

            # 타이틀 + 헤더 접두어
            title_prefix = f"# {sheet_name}\n\n{header_lines}\n"
            prefix_len = len(title_prefix)

            # 데이터 행을 묶어 각 sub-chunk가 임계값 이내가 되도록 분할
            sub_groups: list[list[str]] = []
            current_group: list[str] = []
            current_len = prefix_len

            for row_line in data_lines:
                row_len = len(row_line) + 1  # +1 for '\n'
                if current_group and (current_len + row_len) > _EMBED_TOKEN_LIMIT_CHARS:
                    sub_groups.append(current_group)
                    current_group = [row_line]
                    current_len = prefix_len + row_len
                else:
                    current_group.append(row_line)
                    current_len += row_len

            if current_group:
                sub_groups.append(current_group)

            # sub_groups가 빈 경우 (데이터 행 없음) 단일 청크로 폴백
            if not sub_groups:
                meta = {
                    "source_type": "xlsx",
                    "doc_type": "xlsx",
                    "file_name": file_name,
                    "sheet_name": sheet_name,
                    "chunk_index": chunk_index,
                    "sub_chunk_index": 0,
                    "total_sub_chunks": 1,
                }
                chunks.append({"text": body, "metadata": meta})
                chunk_index += 1
                continue

            total_sub = len(sub_groups)
            for sub_idx, group in enumerate(sub_groups):
                sub_text = title_prefix + "\n".join(group)
                meta = {
                    "source_type": "xlsx",
                    "doc_type": "xlsx",
                    "file_name": file_name,
                    "sheet_name": sheet_name,
                    "chunk_index": chunk_index,
                    "sub_chunk_index": sub_idx,
                    "total_sub_chunks": total_sub,
                }
                chunks.append({"text": sub_text, "metadata": meta})
                chunk_index += 1

    wb.close()
    return chunks


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _shape_text(shape) -> str:
    """pptx shape에서 텍스트 추출 (text_frame, table 셀 모두 처리)."""
    chunks: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                chunks.append(text)
    if shape.has_table:
        for row in shape.table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_cells.append(cell_text)
            if row_cells:
                chunks.append(" | ".join(row_cells))
    return "\n".join(chunks)


def chunk_pptx(file_path: str, stats: dict | None = None) -> list[dict]:
    """PowerPoint 파일: 슬라이드별 1 chunk.

    - python-pptx 사용. title + body shapes의 text + speaker notes 합산
    - 표는 셀별 텍스트 추출
    - bge-m3 한도(``_EMBED_TOKEN_LIMIT_CHARS`` = 6500자) 이내 슬라이드: 단일 청크 (BC 유지)
    - 초과 시: body/notes를 별도 청크로 분리. title은 모든 sub-chunk에 prepend (의미 보존)
      body 단독이 여전히 초과하면 chunk_text로 추가 분할 (warning 로그)
    - metadata: ``{file_name, slide_no, slide_title, doc_type: "pptx", chunk_index, sub_chunk_index, total_sub_chunks}``
      (단일 청크 시에도 sub_chunk_index=0, total_sub_chunks=1 포함 — 일관된 메타 형태)
    - stats: make_ocr_stats() dict를 전달하면 OCR 카운터 + oversize_count 누적 (None이면 생략, BC 유지)
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)

    chunks: list[dict] = []
    chunk_index = 0
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        body_parts: list[str] = []

        # title placeholder 우선 추출
        if slide.shapes.title is not None:
            try:
                slide_title = (slide.shapes.title.text or "").strip()
            except AttributeError:
                slide_title = ""

        for shape in slide.shapes:
            # 타이틀은 위에서 이미 처리 → 본문 부분만
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # alt text는 shape.name(=Picture 3 같은 ID) 아닌 nvPicPr.cNvPr.descr 속성
                try:
                    alt_text = shape._element.nvPicPr.cNvPr.get("descr", "").strip()
                except AttributeError:
                    alt_text = ""
                ocr_text = ""
                try:
                    ocr_text = _ocr_image_blob(shape.image.blob, stats=stats)
                except Exception:
                    pass
                parts = [p for p in (alt_text, ocr_text) if p]
                if parts:
                    body_parts.append(f"[이미지: {' '.join(parts)}]")
                continue
            txt = _shape_text(shape)
            if txt:
                body_parts.append(txt)

        # speaker notes
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()

        # ── 슬라이드 텍스트 조립 ──────────────────────────────────────────
        title_prefix = f"# {slide_title}\n\n" if slide_title else ""
        body_text = "\n".join(body_parts) if body_parts else ""
        notes_section = f"[발표자 노트]\n{notes_text}" if notes_text else ""

        # 완전히 비어있는 슬라이드 건너뜀
        if not body_text and not notes_section and not slide_title:
            continue

        # 전체 합산 텍스트
        sections: list[str] = []
        if slide_title:
            sections.append(f"# {slide_title}")
        if body_text:
            sections.append(body_text)
        if notes_section:
            sections.append(notes_section)

        if not sections:
            continue

        full_body = "\n\n".join(sections)

        def _make_slide_meta(sub_idx: int, total_sub: int) -> dict:
            return {
                "source_type": "pptx",
                "doc_type": "pptx",
                "file_name": file_name,
                "slide_no": slide_no,
                "slide_title": slide_title,
                "chunk_index": chunk_index,
                "sub_chunk_index": sub_idx,
                "total_sub_chunks": total_sub,
            }

        if len(full_body) <= _EMBED_TOKEN_LIMIT_CHARS:
            # ── 단일 청크 (임계값 이내) ────────────────────────────────────
            chunks.append({"text": full_body, "metadata": _make_slide_meta(0, 1)})
            chunk_index += 1
        else:
            # ── 임계값 초과: body/notes 분리 ─────────────────────────────
            if stats is not None:
                stats["oversize_count"] += 1

            sub_texts: list[str] = []

            # body 처리
            if body_text:
                body_chunk_text = (title_prefix + body_text).strip()
                if len(body_chunk_text) <= _EMBED_TOKEN_LIMIT_CHARS:
                    sub_texts.append(body_chunk_text)
                else:
                    # body 단독으로도 초과 — chunk_text로 추가 분할
                    # title_prefix는 모든 조각에 유지해야 하므로
                    # body_text만 분할 대상으로 전달하고 각 조각에 title_prefix를 다시 prepend
                    logger.warning(
                        "chunk_pptx: slide=%d body %d자 > %d 한도 → chunk_text 추가 분할",
                        slide_no, len(body_chunk_text), _EMBED_TOKEN_LIMIT_CHARS,
                    )
                    effective_max = max(500, _EMBED_TOKEN_LIMIT_CHARS - len(title_prefix))
                    sub_pieces = chunk_text(
                        body_text,
                        max_chars=effective_max,
                        overlap=200,
                    )
                    sub_texts.extend(
                        (title_prefix + p["text"]).strip() for p in sub_pieces
                    )

            # notes 처리
            if notes_section:
                notes_chunk_text = (title_prefix + notes_section).strip()
                if len(notes_chunk_text) <= _EMBED_TOKEN_LIMIT_CHARS:
                    sub_texts.append(notes_chunk_text)
                else:
                    logger.warning(
                        "chunk_pptx: slide=%d notes %d자 > %d 한도 → chunk_text 추가 분할",
                        slide_no, len(notes_chunk_text), _EMBED_TOKEN_LIMIT_CHARS,
                    )
                    effective_max = max(500, _EMBED_TOKEN_LIMIT_CHARS - len(title_prefix))
                    sub_pieces = chunk_text(
                        notes_section,
                        max_chars=effective_max,
                        overlap=200,
                    )
                    sub_texts.extend(
                        (title_prefix + p["text"]).strip() for p in sub_pieces
                    )

            # sub_texts 빈 경우 (body도 notes도 없고 title만인 경우) 폴백
            if not sub_texts:
                sub_texts = [full_body]

            total_sub = len(sub_texts)
            for sub_idx, sub_text in enumerate(sub_texts):
                chunks.append({"text": sub_text, "metadata": _make_slide_meta(sub_idx, total_sub)})
                chunk_index += 1

    return chunks
