"""첨부 OCR 처리 — 이미지/PDF/문서 통합 텍스트 추출.

처리 대상:
- 이미지: png/jpg/webp/gif → pytesseract (lang='kor+eng')
- PDF: pdfplumber.extract_text() → 텍스트 없으면 .images blob → _ocr_image_blob 폴백
- docx: python-docx 텍스트 추출
- xlsx: openpyxl 시트별 셀 텍스트 join
- pptx: python-pptx 텍스트 + 도형 이미지 OCR
- txt: 그대로
"""

import io
from pathlib import Path
from typing import Callable
import logging

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# chunking.py:_ocr_image_blob, make_ocr_stats 가져오기
from chunking import _ocr_image_blob, make_ocr_stats

logger = logging.getLogger(__name__)

_IMAGE_MIMES = {"image/png", "image/jpeg", "image/webp", "image/gif"}
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_TEXT_MIME = "text/plain"

_NOOP_CB: Callable[[int], None] = lambda _: None


def extract_text(file_path: Path, mime_type: str) -> str:
    """동기 함수. 첨부 파일에서 텍스트 추출. 실패 시 빈 문자열 반환 + warn 로그.

    반환: 추출된 텍스트 (ocr_stats는 이 함수에서 수집하지 않음)
    """
    text, _ = extract_text_with_stats(file_path, mime_type, _NOOP_CB)
    return text


def extract_text_with_progress(
    file_path: Path, mime_type: str, progress_cb: Callable[[int], None]
) -> str:
    """진행률 콜백 포함 텍스트 추출. progress_cb(0~100)은 동기 컨텍스트(스레드)에서 호출됨.

    OCR 통계가 필요하면 extract_text_with_stats()를 직접 사용할 것.
    """
    text, _ = extract_text_with_stats(file_path, mime_type, progress_cb)
    return text


def extract_text_with_stats(
    file_path: Path,
    mime_type: str,
    progress_cb: Callable[[int], None] = _NOOP_CB,
) -> tuple[str, dict]:
    """텍스트 추출 + OCR 통계 반환.

    Returns:
        (text, ocr_stats) — ocr_stats는 make_ocr_stats() 형태 dict.
        XLSX/DOCX/TXT 처럼 OCR을 하지 않는 포맷은 카운터가 모두 0.
    """
    stats = make_ocr_stats()
    try:
        if mime_type in _IMAGE_MIMES:
            progress_cb(10)
            result = _ocr_image_blob(file_path.read_bytes(), stats=stats)
            progress_cb(100)
            return result, stats
        if mime_type == _PDF_MIME:
            result = _extract_pdf_with_progress(file_path, progress_cb, stats=stats)
            return result, stats
        if mime_type == _DOCX_MIME:
            progress_cb(50)
            result = _extract_docx(file_path)
            progress_cb(100)
            return result, stats
        if mime_type == _XLSX_MIME:
            progress_cb(50)
            result = _extract_xlsx(file_path)
            progress_cb(100)
            return result, stats
        if mime_type == _PPTX_MIME:
            result = _extract_pptx_with_progress(file_path, progress_cb, stats=stats)
            return result, stats
        if mime_type == _TEXT_MIME:
            progress_cb(50)
            result = file_path.read_text(encoding="utf-8", errors="replace")
            progress_cb(100)
            return result, stats
        logger.warning("Unsupported MIME for OCR: %s (%s)", mime_type, file_path)
        return "", stats
    except Exception as exc:
        logger.warning("OCR 실패 (%s, %s): %s", file_path, mime_type, exc)
        return "", stats


def _extract_pdf_with_progress(
    path: Path,
    progress_cb: Callable[[int], None],
    stats: dict | None = None,
) -> str:
    """pdfplumber 우선. 텍스트 미추출 시 페이지 이미지 blob → OCR. 페이지마다 진행률 갱신."""
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        total = max(len(pdf.pages), 1)
        for i, page in enumerate(pdf.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                parts.append(txt)
            else:
                # 텍스트 없는 페이지 — 이미지 OCR 폴백
                for img in page.images:
                    try:
                        img_obj = page.crop((img["x0"], img["top"], img["x1"], img["bottom"])).to_image(resolution=150)
                        buf = io.BytesIO()
                        img_obj.save(buf, format="PNG")
                        parts.append(_ocr_image_blob(buf.getvalue(), stats=stats))
                    except Exception as exc:
                        logger.warning("PDF page OCR fallback fail: %s", exc)
            progress_cb(int((i + 1) / total * 100))
    return "\n\n".join(p for p in parts if p.strip())


def _extract_pptx_with_progress(
    path: Path,
    progress_cb: Callable[[int], None],
    stats: dict | None = None,
) -> str:
    """PPTX 슬라이드별 진행률 갱신 포함 추출."""
    prs = Presentation(path)
    total = max(len(prs.slides), 1)
    parts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {slide_idx}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs)
                    if txt.strip():
                        slide_parts.append(txt)
            if hasattr(shape, "image") and shape.image:
                try:
                    slide_parts.append(_ocr_image_blob(shape.image.blob, stats=stats))
                except Exception:
                    pass
        parts.append("\n".join(slide_parts))
        progress_cb(int(slide_idx / total * 100))
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(path: Path) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        sheet_lines = [f"[Sheet: {sheet.title}]"]
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                sheet_lines.append("\t".join(cells))
        parts.append("\n".join(sheet_lines))
    wb.close()
    return "\n\n".join(parts)
